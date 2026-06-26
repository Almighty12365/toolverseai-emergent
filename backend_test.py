#!/usr/bin/env python3
"""Comprehensive test suite for Toolverse backend endpoints (50+ new endpoints)"""
import requests
import io
import os
from PIL import Image, ImageDraw, ImageFont
import fitz  # PyMuPDF
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt

# Get backend URL from environment
BACKEND_URL = os.getenv("REACT_APP_BACKEND_URL", "https://toolverse-documents.preview.emergentagent.com")
API_BASE = f"{BACKEND_URL}/api"

# Test results tracking
results = {"passed": [], "failed": []}

def log_result(test_name, passed, details=""):
    if passed:
        results["passed"].append(test_name)
        print(f"✓ {test_name}")
    else:
        results["failed"].append({"test": test_name, "details": details})
        print(f"✗ {test_name}: {details}")

def create_test_pdf(pages=3, text_content="Test Page"):
    """Create a test PDF with specified number of pages"""
    doc = fitz.open()
    for i in range(pages):
        page = doc.new_page(width=595, height=842)  # A4
        page.insert_text((50, 100), f"{text_content} {i+1}", fontsize=20)
        page.insert_text((50, 150), f"This is page {i+1} of the test document.", fontsize=12)
        page.insert_text((50, 200), "Lorem ipsum dolor sit amet, consectetur adipiscing elit.", fontsize=10)
    pdf_bytes = doc.tobytes()
    doc.close()
    return pdf_bytes

def create_test_image(width=800, height=600, text="Test Image"):
    """Create a test image"""
    img = Image.new('RGB', (width, height), color=(73, 109, 137))
    d = ImageDraw.Draw(img)
    d.text((width//2-50, height//2), text, fill=(255, 255, 255))
    buf = io.BytesIO()
    img.save(buf, format='PNG')
    return buf.getvalue()

def create_test_image_with_text(text="SAMPLE TEXT"):
    """Create an image with text for OCR testing"""
    img = Image.new('RGB', (400, 200), color='white')
    d = ImageDraw.Draw(img)
    d.text((50, 80), text, fill='black')
    buf = io.BytesIO()
    img.save(buf, format='PNG')
    return buf.getvalue()

def create_test_docx():
    """Create a test DOCX file"""
    doc = Document()
    doc.add_heading('Test Document', 0)
    doc.add_paragraph('This is a test paragraph in the Word document.')
    doc.add_paragraph('Second paragraph with more content.')
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()

def create_test_xlsx():
    """Create a test XLSX file"""
    wb = Workbook()
    ws = wb.active
    ws.title = "Test Sheet"
    ws['A1'] = "Name"
    ws['B1'] = "Value"
    ws['A2'] = "Item 1"
    ws['B2'] = 100
    ws['A3'] = "Item 2"
    ws['B3'] = 200
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()

def create_test_pptx():
    """Create a test PPTX file"""
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "Test Slide 1"
    slide1.placeholders[1].text = "This is the first test slide."
    
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Test Slide 2"
    slide2.placeholders[1].text = "This is the second test slide."
    
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()

def create_test_text():
    """Create test text content"""
    return b"This is a test text file.\n\nIt has multiple paragraphs.\n\nAnd some content for testing."

def create_test_html():
    """Create test HTML content"""
    return b"""<!DOCTYPE html>
<html>
<head><title>Test HTML</title></head>
<body>
<h1>Test Heading</h1>
<p>This is a test paragraph.</p>
<p>Another paragraph with <b>bold</b> text.</p>
</body>
</html>"""

def create_test_csv():
    """Create test CSV content"""
    return b"""Name,Age,City
John Doe,30,New York
Jane Smith,25,Los Angeles
Bob Johnson,35,Chicago"""

print("=" * 80)
print("TOOLVERSE BACKEND API TESTING - 50+ NEW ENDPOINTS")
print("=" * 80)
print(f"Backend URL: {API_BASE}\n")

# ============== CRITICAL TEST: PDF WATERMARK (PREVIOUSLY BROKEN) ==============
print("\n[CRITICAL] Testing PDF Watermark (previously broken)...")
try:
    pdf = create_test_pdf(2, "Watermark Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    data = {"text": "CONFIDENTIAL", "opacity": 0.5, "angle": 45, "color": "#FF0000"}
    r = requests.post(f"{API_BASE}/pdf/watermark", files=files, data=data, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        log_result("PDF Watermark (angle=45)", True)
    else:
        log_result("PDF Watermark (angle=45)", False, f"Status {r.status_code}: {r.text[:200]}")
    
    # Test with valid angles
    for angle in [0, 90, 180, 270]:
        files = {"file": ("test.pdf", pdf, "application/pdf")}
        data = {"text": "TEST", "opacity": 0.3, "angle": angle, "color": "#0000FF"}
        r = requests.post(f"{API_BASE}/pdf/watermark", files=files, data=data, timeout=30)
        if r.status_code == 200:
            log_result(f"PDF Watermark (angle={angle})", True)
        else:
            log_result(f"PDF Watermark (angle={angle})", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Watermark", False, str(e))

# ============== PDF PAGE OPERATIONS ==============
print("\n[PDF PAGE OPERATIONS]")

# 1. Delete Pages
try:
    pdf = create_test_pdf(5, "Delete Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    data = {"pages": "2,4"}
    r = requests.post(f"{API_BASE}/pdf/delete-pages", files=files, data=data, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        # Verify page count
        doc = fitz.open(stream=r.content, filetype="pdf")
        page_count = doc.page_count
        doc.close()
        if page_count == 3:  # 5 - 2 deleted = 3
            log_result("PDF Delete Pages", True)
        else:
            log_result("PDF Delete Pages", False, f"Expected 3 pages, got {page_count}")
    else:
        log_result("PDF Delete Pages", False, f"Status {r.status_code}: {r.text[:200]}")
except Exception as e:
    log_result("PDF Delete Pages", False, str(e))

# 2. Extract Pages
try:
    pdf = create_test_pdf(5, "Extract Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    data = {"pages": "1-2"}
    r = requests.post(f"{API_BASE}/pdf/extract-pages", files=files, data=data, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        doc = fitz.open(stream=r.content, filetype="pdf")
        page_count = doc.page_count
        doc.close()
        if page_count == 2:
            log_result("PDF Extract Pages", True)
        else:
            log_result("PDF Extract Pages", False, f"Expected 2 pages, got {page_count}")
    else:
        log_result("PDF Extract Pages", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Extract Pages", False, str(e))

# 3. Reverse PDF
try:
    pdf = create_test_pdf(3, "Reverse Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    r = requests.post(f"{API_BASE}/pdf/reverse", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        log_result("PDF Reverse", True)
    else:
        log_result("PDF Reverse", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Reverse", False, str(e))

# 4. Duplicate Pages
try:
    pdf = create_test_pdf(3, "Duplicate Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    data = {"pages": "1,3", "count": 2}
    r = requests.post(f"{API_BASE}/pdf/duplicate-pages", files=files, data=data, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        doc = fitz.open(stream=r.content, filetype="pdf")
        page_count = doc.page_count
        doc.close()
        if page_count == 5:  # 3 original + 2 duplicates (pages 1 and 3 duplicated once each)
            log_result("PDF Duplicate Pages", True)
        else:
            log_result("PDF Duplicate Pages", False, f"Expected 5 pages, got {page_count}")
    else:
        log_result("PDF Duplicate Pages", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Duplicate Pages", False, str(e))

# 5. Insert Blank Pages
try:
    pdf = create_test_pdf(2, "Insert Blank Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    data = {"after_page": 1, "count": 2}
    r = requests.post(f"{API_BASE}/pdf/insert-blank", files=files, data=data, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        doc = fitz.open(stream=r.content, filetype="pdf")
        page_count = doc.page_count
        doc.close()
        if page_count == 4:  # 2 + 2 blank
            log_result("PDF Insert Blank Pages", True)
        else:
            log_result("PDF Insert Blank Pages", False, f"Expected 4 pages, got {page_count}")
    else:
        log_result("PDF Insert Blank Pages", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Insert Blank Pages", False, str(e))

# 6. Remove Blank Pages
try:
    # Create PDF with blank page
    doc = fitz.open()
    page1 = doc.new_page(width=595, height=842)
    page1.insert_text((50, 100), "Page with content", fontsize=20)
    doc.new_page(width=595, height=842)  # Blank page
    page3 = doc.new_page(width=595, height=842)
    page3.insert_text((50, 100), "Another page with content", fontsize=20)
    pdf = doc.tobytes()
    doc.close()
    
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    r = requests.post(f"{API_BASE}/pdf/remove-blank", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        doc = fitz.open(stream=r.content, filetype="pdf")
        page_count = doc.page_count
        doc.close()
        if page_count == 2:  # Blank page removed
            log_result("PDF Remove Blank Pages", True)
        else:
            log_result("PDF Remove Blank Pages", False, f"Expected 2 pages, got {page_count}")
    else:
        log_result("PDF Remove Blank Pages", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Remove Blank Pages", False, str(e))

# 7. Crop PDF
try:
    pdf = create_test_pdf(2, "Crop Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    data = {"top": 50, "bottom": 50, "left": 30, "right": 30}
    r = requests.post(f"{API_BASE}/pdf/crop", files=files, data=data, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        log_result("PDF Crop", True)
    else:
        log_result("PDF Crop", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Crop", False, str(e))

# 8. Resize PDF
try:
    pdf = create_test_pdf(2, "Resize Test")
    for size in ["A4", "A3", "A5", "Letter", "Legal"]:
        files = {"file": ("test.pdf", pdf, "application/pdf")}
        data = {"size": size}
        r = requests.post(f"{API_BASE}/pdf/resize", files=files, data=data, timeout=30)
        if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
            log_result(f"PDF Resize ({size})", True)
        else:
            log_result(f"PDF Resize ({size})", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Resize", False, str(e))

# 9. Flatten PDF
try:
    pdf = create_test_pdf(2, "Flatten Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    data = {"dpi": 150}
    r = requests.post(f"{API_BASE}/pdf/flatten", files=files, data=data, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        log_result("PDF Flatten", True)
    else:
        log_result("PDF Flatten", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Flatten", False, str(e))

# 10. Repair PDF
try:
    pdf = create_test_pdf(2, "Repair Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    r = requests.post(f"{API_BASE}/pdf/repair", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        log_result("PDF Repair", True)
    else:
        log_result("PDF Repair", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Repair", False, str(e))

# ============== PDF ANALYSIS ==============
print("\n[PDF ANALYSIS]")

# 11. Extract Images
try:
    # Create PDF with image
    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    img_bytes = create_test_image(200, 200, "Test")
    page.insert_image(fitz.Rect(50, 50, 250, 250), stream=img_bytes)
    pdf = doc.tobytes()
    doc.close()
    
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    r = requests.post(f"{API_BASE}/pdf/extract-images", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/zip":
        log_result("PDF Extract Images", True)
    else:
        log_result("PDF Extract Images", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Extract Images", False, str(e))

# 12. Extract Text
try:
    pdf = create_test_pdf(2, "Extract Text Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    r = requests.post(f"{API_BASE}/pdf/extract-text", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "text/plain":
        text = r.content.decode('utf-8')
        if "Extract Text Test" in text:
            log_result("PDF Extract Text", True)
        else:
            log_result("PDF Extract Text", False, "Text not found in output")
    else:
        log_result("PDF Extract Text", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Extract Text", False, str(e))

# 13. Extract Fonts
try:
    pdf = create_test_pdf(2, "Font Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    r = requests.post(f"{API_BASE}/pdf/extract-fonts", files=files, timeout=30)
    if r.status_code == 200:
        data = r.json()
        if "fonts" in data and "count" in data:
            log_result("PDF Extract Fonts", True)
        else:
            log_result("PDF Extract Fonts", False, "Missing fonts or count in response")
    else:
        log_result("PDF Extract Fonts", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Extract Fonts", False, str(e))

# 14. Search PDF
try:
    pdf = create_test_pdf(2, "Search Test Document")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    data = {"query": "Search Test"}
    r = requests.post(f"{API_BASE}/pdf/search", files=files, data=data, timeout=30)
    if r.status_code == 200:
        result = r.json()
        if "matches" in result and "total" in result and result["total"] > 0:
            log_result("PDF Search", True)
        else:
            log_result("PDF Search", False, f"No matches found: {result}")
    else:
        log_result("PDF Search", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Search", False, str(e))

# 15. Compare PDFs
try:
    pdf1 = create_test_pdf(2, "Document A")
    pdf2 = create_test_pdf(2, "Document B")
    files = [
        ("file1", ("test1.pdf", pdf1, "application/pdf")),
        ("file2", ("test2.pdf", pdf2, "application/pdf"))
    ]
    r = requests.post(f"{API_BASE}/pdf/compare", files=files, timeout=30)
    if r.status_code == 200:
        result = r.json()
        if "fileA_pages" in result and "fileB_pages" in result and "different_pages" in result:
            log_result("PDF Compare", True)
        else:
            log_result("PDF Compare", False, f"Missing fields in response: {result}")
    else:
        log_result("PDF Compare", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Compare", False, str(e))

# 16. Get Metadata
try:
    pdf = create_test_pdf(2, "Metadata Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    r = requests.post(f"{API_BASE}/pdf/metadata", files=files, timeout=30)
    if r.status_code == 200:
        metadata = r.json()
        if "pageCount" in metadata:
            log_result("PDF Get Metadata", True)
        else:
            log_result("PDF Get Metadata", False, "Missing pageCount in metadata")
    else:
        log_result("PDF Get Metadata", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Get Metadata", False, str(e))

# 17. Set Metadata
try:
    pdf = create_test_pdf(2, "Set Metadata Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    data = {"title": "Test Title", "author": "Test Author", "subject": "Test Subject", "keywords": "test, keywords"}
    r = requests.post(f"{API_BASE}/pdf/set-metadata", files=files, data=data, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        log_result("PDF Set Metadata", True)
    else:
        log_result("PDF Set Metadata", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Set Metadata", False, str(e))

# 18. Remove Metadata
try:
    pdf = create_test_pdf(2, "Remove Metadata Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    r = requests.post(f"{API_BASE}/pdf/remove-metadata", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        log_result("PDF Remove Metadata", True)
    else:
        log_result("PDF Remove Metadata", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Remove Metadata", False, str(e))

# ============== PDF EDITING ==============
print("\n[PDF EDITING]")

# 19. Header/Footer
try:
    pdf = create_test_pdf(2, "Header Footer Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    data = {"header": "Test Header", "footer": "Test Footer"}
    r = requests.post(f"{API_BASE}/pdf/header-footer", files=files, data=data, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        log_result("PDF Header/Footer", True)
    else:
        log_result("PDF Header/Footer", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Header/Footer", False, str(e))

# 20. Bates Numbering
try:
    pdf = create_test_pdf(3, "Bates Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    data = {"prefix": "DOC-", "start": 1, "digits": 6}
    r = requests.post(f"{API_BASE}/pdf/bates", files=files, data=data, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        log_result("PDF Bates Numbering", True)
    else:
        log_result("PDF Bates Numbering", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF Bates Numbering", False, str(e))

# ============== CONVERSIONS ==============
print("\n[CONVERSIONS]")

# 21. PDF to TXT
try:
    pdf = create_test_pdf(2, "PDF to TXT Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    r = requests.post(f"{API_BASE}/pdf/to-txt", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "text/plain":
        log_result("PDF to TXT", True)
    else:
        log_result("PDF to TXT", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF to TXT", False, str(e))

# 22. TXT to PDF
try:
    txt = create_test_text()
    files = {"file": ("test.txt", txt, "text/plain")}
    r = requests.post(f"{API_BASE}/txt/to-pdf", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        log_result("TXT to PDF", True)
    else:
        log_result("TXT to PDF", False, f"Status {r.status_code}")
except Exception as e:
    log_result("TXT to PDF", False, str(e))

# 23. PDF to HTML
try:
    pdf = create_test_pdf(2, "PDF to HTML Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    r = requests.post(f"{API_BASE}/pdf/to-html", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "text/html":
        log_result("PDF to HTML", True)
    else:
        log_result("PDF to HTML", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF to HTML", False, str(e))

# 24. HTML to PDF
try:
    html = create_test_html()
    files = {"file": ("test.html", html, "text/html")}
    r = requests.post(f"{API_BASE}/html/to-pdf", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        log_result("HTML to PDF", True)
    else:
        log_result("HTML to PDF", False, f"Status {r.status_code}")
except Exception as e:
    log_result("HTML to PDF", False, str(e))

# 25. PDF to CSV
try:
    pdf = create_test_pdf(2, "PDF to CSV Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    r = requests.post(f"{API_BASE}/pdf/to-csv", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "text/csv":
        log_result("PDF to CSV", True)
    else:
        log_result("PDF to CSV", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF to CSV", False, str(e))

# 26. CSV to PDF
try:
    csv = create_test_csv()
    files = {"file": ("test.csv", csv, "text/csv")}
    r = requests.post(f"{API_BASE}/csv/to-pdf", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        log_result("CSV to PDF", True)
    else:
        log_result("CSV to PDF", False, f"Status {r.status_code}")
except Exception as e:
    log_result("CSV to PDF", False, str(e))

# 27. PDF to SVG
try:
    pdf = create_test_pdf(2, "PDF to SVG Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    r = requests.post(f"{API_BASE}/pdf/to-svg", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/zip":
        log_result("PDF to SVG", True)
    else:
        log_result("PDF to SVG", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF to SVG", False, str(e))

# 28. PDF to Word
try:
    pdf = create_test_pdf(2, "PDF to Word Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    r = requests.post(f"{API_BASE}/pdf/to-word", files=files, timeout=30)
    if r.status_code == 200 and "wordprocessingml" in r.headers.get("content-type", ""):
        log_result("PDF to Word", True)
    else:
        log_result("PDF to Word", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF to Word", False, str(e))

# 29. Word to PDF
try:
    docx = create_test_docx()
    files = {"file": ("test.docx", docx, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
    r = requests.post(f"{API_BASE}/word/to-pdf", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        log_result("Word to PDF", True)
    else:
        log_result("Word to PDF", False, f"Status {r.status_code}")
except Exception as e:
    log_result("Word to PDF", False, str(e))

# 30. PDF to Excel
try:
    pdf = create_test_pdf(2, "PDF to Excel Test")
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    r = requests.post(f"{API_BASE}/pdf/to-excel", files=files, timeout=30)
    if r.status_code == 200 and "spreadsheetml" in r.headers.get("content-type", ""):
        log_result("PDF to Excel", True)
    else:
        log_result("PDF to Excel", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PDF to Excel", False, str(e))

# 31. Excel to PDF
try:
    xlsx = create_test_xlsx()
    files = {"file": ("test.xlsx", xlsx, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")}
    r = requests.post(f"{API_BASE}/excel/to-pdf", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        log_result("Excel to PDF", True)
    else:
        log_result("Excel to PDF", False, f"Status {r.status_code}")
except Exception as e:
    log_result("Excel to PDF", False, str(e))

# 32. PPT to PDF
try:
    pptx = create_test_pptx()
    files = {"file": ("test.pptx", pptx, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
    r = requests.post(f"{API_BASE}/ppt/to-pdf", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/pdf":
        log_result("PPT to PDF", True)
    else:
        log_result("PPT to PDF", False, f"Status {r.status_code}")
except Exception as e:
    log_result("PPT to PDF", False, str(e))

# ============== OCR ==============
print("\n[OCR]")

# 33. OCR PDF
try:
    # Create PDF with text image
    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    img_bytes = create_test_image_with_text("HELLO WORLD")
    page.insert_image(fitz.Rect(50, 50, 450, 250), stream=img_bytes)
    pdf = doc.tobytes()
    doc.close()
    
    files = {"file": ("test.pdf", pdf, "application/pdf")}
    data = {"lang": "eng"}
    r = requests.post(f"{API_BASE}/ocr/pdf", files=files, data=data, timeout=60)
    if r.status_code == 200:
        result = r.json()
        if "text" in result and "pages" in result:
            log_result("OCR PDF", True)
        else:
            log_result("OCR PDF", False, f"Missing fields: {result}")
    else:
        log_result("OCR PDF", False, f"Status {r.status_code}: {r.text[:200]}")
except Exception as e:
    log_result("OCR PDF", False, str(e))

# 34. OCR Image
try:
    img = create_test_image_with_text("TEST OCR IMAGE")
    files = {"file": ("test.png", img, "image/png")}
    data = {"lang": "eng"}
    r = requests.post(f"{API_BASE}/ocr/image", files=files, data=data, timeout=60)
    if r.status_code == 200:
        result = r.json()
        if "text" in result:
            log_result("OCR Image", True)
        else:
            log_result("OCR Image", False, "Missing text field")
    else:
        log_result("OCR Image", False, f"Status {r.status_code}")
except Exception as e:
    log_result("OCR Image", False, str(e))

# ============== IMAGE EXTRAS ==============
print("\n[IMAGE EXTRAS]")

# 35. Image Crop
try:
    img = create_test_image(800, 600)
    files = {"file": ("test.png", img, "image/png")}
    data = {"left": 100, "top": 100, "right": 100, "bottom": 100}
    r = requests.post(f"{API_BASE}/image/crop", files=files, data=data, timeout=30)
    if r.status_code == 200 and "image" in r.headers.get("content-type", ""):
        log_result("Image Crop", True)
    else:
        log_result("Image Crop", False, f"Status {r.status_code}")
except Exception as e:
    log_result("Image Crop", False, str(e))

# 36. Image Rotate
try:
    img = create_test_image(400, 300)
    files = {"file": ("test.png", img, "image/png")}
    data = {"angle": 45}
    r = requests.post(f"{API_BASE}/image/rotate", files=files, data=data, timeout=30)
    if r.status_code == 200 and "image" in r.headers.get("content-type", ""):
        log_result("Image Rotate", True)
    else:
        log_result("Image Rotate", False, f"Status {r.status_code}")
except Exception as e:
    log_result("Image Rotate", False, str(e))

# 37. Image Flip
try:
    img = create_test_image(400, 300)
    for direction in ["horizontal", "vertical"]:
        files = {"file": ("test.png", img, "image/png")}
        data = {"direction": direction}
        r = requests.post(f"{API_BASE}/image/flip", files=files, data=data, timeout=30)
        if r.status_code == 200 and "image" in r.headers.get("content-type", ""):
            log_result(f"Image Flip ({direction})", True)
        else:
            log_result(f"Image Flip ({direction})", False, f"Status {r.status_code}")
except Exception as e:
    log_result("Image Flip", False, str(e))

# 38. Image Border
try:
    img = create_test_image(400, 300)
    files = {"file": ("test.png", img, "image/png")}
    data = {"thickness": 20, "color": "#FF0000"}
    r = requests.post(f"{API_BASE}/image/border", files=files, data=data, timeout=30)
    if r.status_code == 200 and "image" in r.headers.get("content-type", ""):
        log_result("Image Border", True)
    else:
        log_result("Image Border", False, f"Status {r.status_code}")
except Exception as e:
    log_result("Image Border", False, str(e))

# 39. Image Round Corners
try:
    img = create_test_image(400, 300)
    files = {"file": ("test.png", img, "image/png")}
    data = {"radius": 40}
    r = requests.post(f"{API_BASE}/image/round-corners", files=files, data=data, timeout=30)
    if r.status_code == 200 and "image" in r.headers.get("content-type", ""):
        log_result("Image Round Corners", True)
    else:
        log_result("Image Round Corners", False, f"Status {r.status_code}")
except Exception as e:
    log_result("Image Round Corners", False, str(e))

# 40. Image DPI
try:
    img = create_test_image(400, 300)
    files = {"file": ("test.png", img, "image/png")}
    data = {"dpi": 300}
    r = requests.post(f"{API_BASE}/image/dpi", files=files, data=data, timeout=30)
    if r.status_code == 200 and "image" in r.headers.get("content-type", ""):
        log_result("Image DPI", True)
    else:
        log_result("Image DPI", False, f"Status {r.status_code}")
except Exception as e:
    log_result("Image DPI", False, str(e))

# 41. Image Metadata
try:
    img = create_test_image(400, 300)
    files = {"file": ("test.png", img, "image/png")}
    r = requests.post(f"{API_BASE}/image/metadata", files=files, timeout=30)
    if r.status_code == 200:
        metadata = r.json()
        if "format" in metadata and "size" in metadata:
            log_result("Image Metadata", True)
        else:
            log_result("Image Metadata", False, "Missing fields in metadata")
    else:
        log_result("Image Metadata", False, f"Status {r.status_code}")
except Exception as e:
    log_result("Image Metadata", False, str(e))

# 42. Image Remove Metadata
try:
    img = create_test_image(400, 300)
    files = {"file": ("test.png", img, "image/png")}
    r = requests.post(f"{API_BASE}/image/remove-metadata", files=files, timeout=30)
    if r.status_code == 200 and "image" in r.headers.get("content-type", ""):
        log_result("Image Remove Metadata", True)
    else:
        log_result("Image Remove Metadata", False, f"Status {r.status_code}")
except Exception as e:
    log_result("Image Remove Metadata", False, str(e))

# 43. Image Contact Sheet
try:
    imgs = [create_test_image(200, 200, f"Img{i}") for i in range(6)]
    files = [("files", (f"test{i}.png", img, "image/png")) for i, img in enumerate(imgs)]
    data = {"cols": 3, "thumb": 150}
    r = requests.post(f"{API_BASE}/image/contact-sheet", files=files, data=data, timeout=30)
    if r.status_code == 200 and "image" in r.headers.get("content-type", ""):
        log_result("Image Contact Sheet", True)
    else:
        log_result("Image Contact Sheet", False, f"Status {r.status_code}")
except Exception as e:
    log_result("Image Contact Sheet", False, str(e))

# 44. Image Favicon
try:
    img = create_test_image(256, 256, "ICON")
    files = {"file": ("test.png", img, "image/png")}
    r = requests.post(f"{API_BASE}/image/favicon", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/zip":
        log_result("Image Favicon", True)
    else:
        log_result("Image Favicon", False, f"Status {r.status_code}")
except Exception as e:
    log_result("Image Favicon", False, str(e))

# ============== BARCODE ==============
print("\n[BARCODE]")

# 45. Barcode Generate
try:
    for kind in ["code128", "code39", "ean13", "ean8", "upca"]:
        # Generate appropriate data for each barcode type
        if kind == "ean13":
            data_val = "5901234123457"
        elif kind == "ean8":
            data_val = "12345670"
        elif kind == "upca":
            data_val = "123456789012"
        else:
            data_val = "TEST123"
        
        data = {"data": data_val, "kind": kind}
        r = requests.post(f"{API_BASE}/barcode/generate", data=data, timeout=30)
        if r.status_code == 200 and "image" in r.headers.get("content-type", ""):
            log_result(f"Barcode Generate ({kind})", True)
        else:
            log_result(f"Barcode Generate ({kind})", False, f"Status {r.status_code}: {r.text[:200]}")
except Exception as e:
    log_result("Barcode Generate", False, str(e))

# ============== ARCHIVES ==============
print("\n[ARCHIVES]")

# 46. ZIP Create
try:
    files = [
        ("files", ("file1.txt", b"Content of file 1", "text/plain")),
        ("files", ("file2.txt", b"Content of file 2", "text/plain")),
        ("files", ("file3.txt", b"Content of file 3", "text/plain"))
    ]
    r = requests.post(f"{API_BASE}/zip/create", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/zip":
        log_result("ZIP Create", True)
    else:
        log_result("ZIP Create", False, f"Status {r.status_code}")
except Exception as e:
    log_result("ZIP Create", False, str(e))

# 47. ZIP Extract
try:
    # First create a zip
    files = [
        ("files", ("file1.txt", b"Content 1", "text/plain")),
        ("files", ("file2.txt", b"Content 2", "text/plain"))
    ]
    r = requests.post(f"{API_BASE}/zip/create", files=files, timeout=30)
    if r.status_code == 200:
        zip_data = r.content
        files = {"file": ("test.zip", zip_data, "application/zip")}
        r = requests.post(f"{API_BASE}/zip/extract", files=files, timeout=30)
        if r.status_code == 200:
            result = r.json()
            if "files" in result and "count" in result and result["count"] == 2:
                log_result("ZIP Extract", True)
            else:
                log_result("ZIP Extract", False, f"Unexpected result: {result}")
        else:
            log_result("ZIP Extract", False, f"Status {r.status_code}")
    else:
        log_result("ZIP Extract", False, "Failed to create test zip")
except Exception as e:
    log_result("ZIP Extract", False, str(e))

# 48. TAR Create
try:
    files = [
        ("files", ("file1.txt", b"Content of file 1", "text/plain")),
        ("files", ("file2.txt", b"Content of file 2", "text/plain"))
    ]
    r = requests.post(f"{API_BASE}/tar/create", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/x-tar":
        log_result("TAR Create", True)
    else:
        log_result("TAR Create", False, f"Status {r.status_code}")
except Exception as e:
    log_result("TAR Create", False, str(e))

# 49. GZIP Compress
try:
    txt = create_test_text()
    files = {"file": ("test.txt", txt, "text/plain")}
    r = requests.post(f"{API_BASE}/gzip/compress", files=files, timeout=30)
    if r.status_code == 200 and r.headers.get("content-type") == "application/gzip":
        log_result("GZIP Compress", True)
    else:
        log_result("GZIP Compress", False, f"Status {r.status_code}")
except Exception as e:
    log_result("GZIP Compress", False, str(e))

# 50. GZIP Decompress
try:
    txt = create_test_text()
    files = {"file": ("test.txt", txt, "text/plain")}
    r = requests.post(f"{API_BASE}/gzip/compress", files=files, timeout=30)
    if r.status_code == 200:
        gz_data = r.content
        files = {"file": ("test.txt.gz", gz_data, "application/gzip")}
        r = requests.post(f"{API_BASE}/gzip/decompress", files=files, timeout=30)
        if r.status_code == 200:
            log_result("GZIP Decompress", True)
        else:
            log_result("GZIP Decompress", False, f"Status {r.status_code}")
    else:
        log_result("GZIP Decompress", False, "Failed to create test gzip")
except Exception as e:
    log_result("GZIP Decompress", False, str(e))

# ============== SUMMARY ==============
print("\n" + "=" * 80)
print("TEST SUMMARY")
print("=" * 80)
print(f"✓ PASSED: {len(results['passed'])} tests")
print(f"✗ FAILED: {len(results['failed'])} tests")

if results["failed"]:
    print("\nFAILED TESTS:")
    for fail in results["failed"]:
        print(f"  ✗ {fail['test']}")
        print(f"    Details: {fail['details']}")

print("\n" + "=" * 80)
print(f"Total tests: {len(results['passed']) + len(results['failed'])}")
print("=" * 80)
