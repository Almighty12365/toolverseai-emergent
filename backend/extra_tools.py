"""Additional Toolverse endpoints — PDF, image, conversion, archive, OCR, etc."""
from fastapi import APIRouter, UploadFile, File, Form, HTTPException
from fastapi.responses import StreamingResponse
import io, os, zipfile, tarfile, gzip, uuid, hashlib, secrets, string, json, re, csv as csvmod
from typing import List, Optional
import fitz  # PyMuPDF
from PIL import Image, ImageOps, ImageDraw
import pytesseract
from docx import Document
from openpyxl import load_workbook, Workbook
from pptx import Presentation
from reportlab.pdfgen import canvas as rl_canvas
from reportlab.lib.pagesizes import A4, letter
from reportlab.lib.units import mm
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from xhtml2pdf import pisa
import barcode
from barcode.writer import ImageWriter

router = APIRouter(prefix="/api")


def _stream(data: bytes, filename: str, media_type: str = "application/pdf"):
    return StreamingResponse(
        io.BytesIO(data),
        media_type=media_type,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )


# ============== PDF PAGE OPERATIONS ==============

@router.post("/pdf/delete-pages")
async def delete_pages(file: UploadFile = File(...), pages: str = Form(...)):
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        to_del = set()
        for ch in pages.split(","):
            ch = ch.strip()
            if "-" in ch:
                a, b = ch.split("-")
                to_del.update(range(int(a) - 1, int(b)))
            else:
                to_del.add(int(ch) - 1)
        new = fitz.open()
        for i in range(doc.page_count):
            if i not in to_del:
                new.insert_pdf(doc, from_page=i, to_page=i)
        out = new.tobytes()
        new.close(); doc.close()
        return _stream(out, "pages-deleted.pdf")
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


@router.post("/pdf/extract-pages")
async def extract_pages(file: UploadFile = File(...), pages: str = Form(...)):
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        keep = []
        for ch in pages.split(","):
            ch = ch.strip()
            if "-" in ch:
                a, b = ch.split("-"); keep.extend(range(int(a) - 1, int(b)))
            else:
                keep.append(int(ch) - 1)
        new = fitz.open()
        for i in keep:
            if 0 <= i < doc.page_count:
                new.insert_pdf(doc, from_page=i, to_page=i)
        out = new.tobytes(); new.close(); doc.close()
        return _stream(out, "extracted.pdf")
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


@router.post("/pdf/reverse")
async def reverse_pdf(file: UploadFile = File(...)):
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        new = fitz.open()
        for i in reversed(range(doc.page_count)):
            new.insert_pdf(doc, from_page=i, to_page=i)
        out = new.tobytes(); new.close(); doc.close()
        return _stream(out, "reversed.pdf")
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


@router.post("/pdf/duplicate-pages")
async def duplicate_pages(file: UploadFile = File(...), pages: str = Form(...), count: int = Form(2)):
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        dup = set()
        for ch in pages.split(","):
            ch = ch.strip()
            if "-" in ch:
                a, b = ch.split("-"); dup.update(range(int(a) - 1, int(b)))
            else:
                dup.add(int(ch) - 1)
        new = fitz.open()
        for i in range(doc.page_count):
            new.insert_pdf(doc, from_page=i, to_page=i)
            if i in dup:
                for _ in range(count - 1):
                    new.insert_pdf(doc, from_page=i, to_page=i)
        out = new.tobytes(); new.close(); doc.close()
        return _stream(out, "duplicated.pdf")
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


@router.post("/pdf/insert-blank")
async def insert_blank(file: UploadFile = File(...), after_page: int = Form(...), count: int = Form(1)):
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        idx = max(0, min(after_page, doc.page_count))
        ref = doc[max(0, idx - 1)] if doc.page_count else None
        w, h = (ref.rect.width, ref.rect.height) if ref else (595, 842)
        for _ in range(count):
            doc.insert_page(idx, width=w, height=h)
            idx += 1
        out = doc.tobytes(); doc.close()
        return _stream(out, "with-blanks.pdf")
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


@router.post("/pdf/remove-blank")
async def remove_blank(file: UploadFile = File(...)):
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        new = fitz.open()
        for i, page in enumerate(doc):
            text = page.get_text().strip()
            images = page.get_images()
            if text or images:
                new.insert_pdf(doc, from_page=i, to_page=i)
        out = new.tobytes(); new.close(); doc.close()
        return _stream(out, "blanks-removed.pdf")
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


@router.post("/pdf/crop")
async def crop_pdf(file: UploadFile = File(...), top: float = Form(0), bottom: float = Form(0), left: float = Form(0), right: float = Form(0)):
    """Crop margins (in points)."""
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        for page in doc:
            r = page.rect
            new_rect = fitz.Rect(r.x0 + left, r.y0 + top, r.x1 - right, r.y1 - bottom)
            page.set_cropbox(new_rect)
        out = doc.tobytes(); doc.close()
        return _stream(out, "cropped.pdf")
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


@router.post("/pdf/resize")
async def resize_pdf(file: UploadFile = File(...), size: str = Form("A4")):
    """Resize all pages to specified size (A4/Letter/A3/A5)."""
    data = await file.read()
    src = fitz.open(stream=data, filetype="pdf")
    sizes = {
        "A4": (595, 842), "A3": (842, 1191), "A5": (420, 595),
        "Letter": (612, 792), "Legal": (612, 1008),
    }
    tw, th = sizes.get(size, (595, 842))
    try:
        new = fitz.open()
        for page in src:
            new_page = new.new_page(width=tw, height=th)
            new_page.show_pdf_page(new_page.rect, src, page.number)
        out = new.tobytes(); new.close(); src.close()
        return _stream(out, "resized.pdf")
    except Exception as e:
        src.close(); raise HTTPException(500, str(e))


@router.post("/pdf/flatten")
async def flatten_pdf(file: UploadFile = File(...), dpi: int = Form(150)):
    data = await file.read()
    src = fitz.open(stream=data, filetype="pdf")
    try:
        new = fitz.open()
        mat = fitz.Matrix(dpi / 72, dpi / 72)
        for page in src:
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes("png")
            r = page.rect
            p = new.new_page(width=r.width, height=r.height)
            p.insert_image(p.rect, stream=img_bytes)
        out = new.tobytes(); new.close(); src.close()
        return _stream(out, "flattened.pdf")
    except Exception as e:
        src.close(); raise HTTPException(500, str(e))


@router.post("/pdf/repair")
async def repair_pdf(file: UploadFile = File(...)):
    data = await file.read()
    try:
        doc = fitz.open(stream=data, filetype="pdf")
        out = doc.tobytes(garbage=4, clean=True, deflate=True)
        doc.close()
        return _stream(out, "repaired.pdf")
    except Exception as e:
        raise HTTPException(500, f"Repair failed: {str(e)}")


@router.post("/pdf/extract-images")
async def extract_images(file: UploadFile = File(...)):
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        zbuf = io.BytesIO()
        with zipfile.ZipFile(zbuf, "w", zipfile.ZIP_DEFLATED) as zf:
            count = 0
            for pno, page in enumerate(doc, 1):
                for img in page.get_images(full=True):
                    xref = img[0]
                    try:
                        info = doc.extract_image(xref)
                        ext = info.get("ext", "png")
                        zf.writestr(f"page{pno}-img{count}.{ext}", info["image"])
                        count += 1
                    except Exception:
                        continue
        doc.close(); zbuf.seek(0)
        return _stream(zbuf.read(), "extracted-images.zip", "application/zip")
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


@router.post("/pdf/extract-text")
async def extract_text(file: UploadFile = File(...)):
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        text = "".join(page.get_text() + "\n\n" for page in doc)
        doc.close()
        return _stream(text.encode("utf-8"), "extracted.txt", "text/plain")
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


@router.post("/pdf/extract-fonts")
async def extract_fonts(file: UploadFile = File(...)):
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        fonts = set()
        for page in doc:
            for f in page.get_fonts(full=True):
                fonts.add(f[3])
        doc.close()
        return {"fonts": sorted(fonts), "count": len(fonts)}
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


@router.post("/pdf/search")
async def search_pdf(file: UploadFile = File(...), query: str = Form(...)):
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        hits = []
        for pno, page in enumerate(doc, 1):
            rects = page.search_for(query)
            for r in rects:
                hits.append({"page": pno, "x": r.x0, "y": r.y0, "w": r.width, "h": r.height})
        doc.close()
        return {"matches": hits, "total": len(hits)}
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


@router.post("/pdf/compare")
async def compare_pdfs(file1: UploadFile = File(...), file2: UploadFile = File(...)):
    d1 = await file1.read(); d2 = await file2.read()
    a = fitz.open(stream=d1, filetype="pdf"); b = fitz.open(stream=d2, filetype="pdf")
    try:
        diff_pages = []
        for i in range(max(a.page_count, b.page_count)):
            ta = a[i].get_text() if i < a.page_count else ""
            tb = b[i].get_text() if i < b.page_count else ""
            if ta != tb:
                diff_pages.append(i + 1)
        result = {
            "fileA_pages": a.page_count, "fileB_pages": b.page_count,
            "different_pages": diff_pages, "identical": len(diff_pages) == 0
        }
        a.close(); b.close()
        return result
    except Exception as e:
        a.close(); b.close(); raise HTTPException(500, str(e))


# ============== METADATA / PERMISSIONS ==============

@router.post("/pdf/metadata")
async def get_metadata(file: UploadFile = File(...)):
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        m = dict(doc.metadata or {})
        m["pageCount"] = doc.page_count
        doc.close()
        return m
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


@router.post("/pdf/set-metadata")
async def set_metadata(
    file: UploadFile = File(...),
    title: str = Form(""), author: str = Form(""),
    subject: str = Form(""), keywords: str = Form("")
):
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        meta = doc.metadata or {}
        if title: meta["title"] = title
        if author: meta["author"] = author
        if subject: meta["subject"] = subject
        if keywords: meta["keywords"] = keywords
        doc.set_metadata(meta)
        out = doc.tobytes(); doc.close()
        return _stream(out, "with-metadata.pdf")
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


@router.post("/pdf/remove-metadata")
async def remove_metadata(file: UploadFile = File(...)):
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        doc.set_metadata({})
        out = doc.tobytes(); doc.close()
        return _stream(out, "stripped.pdf")
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


# ============== HEADER / FOOTER / BATES ==============

@router.post("/pdf/header-footer")
async def header_footer(
    file: UploadFile = File(...),
    header: str = Form(""), footer: str = Form(""),
):
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        for page in doc:
            r = page.rect
            if header:
                page.insert_text((30, 25), header, fontsize=10, color=(0.2, 0.2, 0.2))
            if footer:
                page.insert_text((30, r.height - 20), footer, fontsize=10, color=(0.2, 0.2, 0.2))
        out = doc.tobytes(); doc.close()
        return _stream(out, "header-footer.pdf")
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


@router.post("/pdf/bates")
async def bates_numbering(
    file: UploadFile = File(...),
    prefix: str = Form(""), start: int = Form(1), digits: int = Form(6)
):
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        for i, page in enumerate(doc):
            label = f"{prefix}{str(start + i).zfill(digits)}"
            r = page.rect
            tw = fitz.get_text_length(label, fontname="helv", fontsize=10)
            page.insert_text((r.width - tw - 30, r.height - 20), label, fontsize=10, color=(0.2, 0.2, 0.2))
        out = doc.tobytes(); doc.close()
        return _stream(out, "bates.pdf")
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


# ============== CONVERSIONS ==============

@router.post("/pdf/to-txt")
async def pdf_to_txt(file: UploadFile = File(...)):
    return await extract_text(file)


@router.post("/txt/to-pdf")
async def txt_to_pdf(file: UploadFile = File(...)):
    data = await file.read()
    text = data.decode("utf-8", errors="ignore")
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4)
    styles = getSampleStyleSheet()
    story = []
    for para in text.split("\n\n"):
        story.append(Paragraph(para.replace("\n", "<br/>") or "&nbsp;", styles["BodyText"]))
        story.append(Spacer(1, 6))
    doc.build(story)
    return _stream(buf.getvalue(), "from-text.pdf")


@router.post("/pdf/to-html")
async def pdf_to_html(file: UploadFile = File(...)):
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        html = "<html><head><meta charset='utf-8'><title>Converted</title></head><body>"
        for page in doc:
            html += page.get_text("html") + "<hr/>"
        html += "</body></html>"
        doc.close()
        return _stream(html.encode("utf-8"), "converted.html", "text/html")
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


@router.post("/html/to-pdf")
async def html_to_pdf(file: UploadFile = File(...)):
    data = await file.read()
    html = data.decode("utf-8", errors="ignore")
    out = io.BytesIO()
    pisa.CreatePDF(html, dest=out)
    return _stream(out.getvalue(), "from-html.pdf")


@router.post("/pdf/to-csv")
async def pdf_to_csv(file: UploadFile = File(...)):
    """Extract text lines as CSV rows (basic)."""
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        out = io.StringIO(); w = csvmod.writer(out)
        w.writerow(["page", "text"])
        for i, page in enumerate(doc, 1):
            for line in page.get_text().splitlines():
                if line.strip():
                    w.writerow([i, line])
        doc.close()
        return _stream(out.getvalue().encode("utf-8"), "out.csv", "text/csv")
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


@router.post("/csv/to-pdf")
async def csv_to_pdf(file: UploadFile = File(...)):
    data = await file.read()
    text = data.decode("utf-8", errors="ignore")
    reader = csvmod.reader(io.StringIO(text))
    buf = io.BytesIO()
    c = rl_canvas.Canvas(buf, pagesize=A4)
    width, height = A4
    y = height - 40
    for row in reader:
        line = " | ".join(row)
        c.drawString(30, y, line[:120])
        y -= 14
        if y < 40:
            c.showPage(); y = height - 40
    c.save()
    return _stream(buf.getvalue(), "from-csv.pdf")


@router.post("/pdf/to-word")
async def pdf_to_word(file: UploadFile = File(...)):
    """Basic PDF -> DOCX (text only)."""
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        d = Document()
        for page in doc:
            for para in page.get_text().split("\n"):
                if para.strip():
                    d.add_paragraph(para)
            d.add_page_break()
        doc.close()
        out = io.BytesIO(); d.save(out)
        return _stream(out.getvalue(), "converted.docx",
                       "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


@router.post("/word/to-pdf")
async def word_to_pdf(file: UploadFile = File(...)):
    data = await file.read()
    bio = io.BytesIO(data)
    try:
        d = Document(bio)
        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        for p in d.paragraphs:
            txt = (p.text or "").strip()
            story.append(Paragraph(txt or "&nbsp;", styles["BodyText"]))
        doc.build(story)
        return _stream(buf.getvalue(), "from-word.pdf")
    except Exception as e:
        raise HTTPException(500, f"Word→PDF failed: {str(e)}")


@router.post("/pdf/to-excel")
async def pdf_to_excel(file: UploadFile = File(...)):
    """Extract each page's text into a new worksheet row-by-row."""
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        wb = Workbook(); wb.remove(wb.active)
        for i, page in enumerate(doc, 1):
            ws = wb.create_sheet(title=f"Page {i}")
            for line in page.get_text().splitlines():
                ws.append([line])
        doc.close()
        out = io.BytesIO(); wb.save(out)
        return _stream(out.getvalue(), "converted.xlsx",
                       "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


@router.post("/excel/to-pdf")
async def excel_to_pdf(file: UploadFile = File(...)):
    data = await file.read()
    wb = load_workbook(io.BytesIO(data), data_only=True)
    buf = io.BytesIO()
    c = rl_canvas.Canvas(buf, pagesize=A4)
    width, height = A4
    for sname in wb.sheetnames:
        ws = wb[sname]
        c.setFont("Helvetica-Bold", 12); c.drawString(30, height - 40, f"Sheet: {sname}")
        c.setFont("Helvetica", 9)
        y = height - 60
        for row in ws.iter_rows(values_only=True):
            line = " | ".join("" if v is None else str(v) for v in row)
            c.drawString(30, y, line[:140])
            y -= 12
            if y < 40:
                c.showPage(); y = height - 40
                c.setFont("Helvetica", 9)
        c.showPage()
    c.save()
    return _stream(buf.getvalue(), "from-excel.pdf")


@router.post("/ppt/to-pdf")
async def ppt_to_pdf(file: UploadFile = File(...)):
    data = await file.read()
    prs = Presentation(io.BytesIO(data))
    buf = io.BytesIO()
    c = rl_canvas.Canvas(buf, pagesize=letter)
    width, height = letter
    for idx, slide in enumerate(prs.slides, 1):
        c.setFont("Helvetica-Bold", 14); c.drawString(30, height - 50, f"Slide {idx}")
        c.setFont("Helvetica", 11)
        y = height - 80
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                for line in shape.text.splitlines():
                    c.drawString(40, y, line[:110]); y -= 14
                    if y < 40:
                        c.showPage(); y = height - 50
        c.showPage()
    c.save()
    return _stream(buf.getvalue(), "from-ppt.pdf")


@router.post("/pdf/to-svg")
async def pdf_to_svg(file: UploadFile = File(...)):
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        zbuf = io.BytesIO()
        with zipfile.ZipFile(zbuf, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, page in enumerate(doc, 1):
                svg = page.get_svg_image()
                zf.writestr(f"page-{i}.svg", svg)
        doc.close(); zbuf.seek(0)
        return _stream(zbuf.read(), "svg-pages.zip", "application/zip")
    except Exception as e:
        doc.close(); raise HTTPException(500, str(e))


# ============== OCR ==============

@router.post("/ocr/pdf")
async def ocr_pdf(file: UploadFile = File(...), lang: str = Form("eng")):
    data = await file.read()
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        full = []
        for page in doc:
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            full.append(pytesseract.image_to_string(img, lang=lang))
        doc.close()
        return {"text": "\n\n".join(full), "pages": len(full)}
    except Exception as e:
        doc.close(); raise HTTPException(500, f"OCR failed: {str(e)}")


@router.post("/ocr/image")
async def ocr_image(file: UploadFile = File(...), lang: str = Form("eng")):
    data = await file.read()
    img = Image.open(io.BytesIO(data))
    try:
        text = pytesseract.image_to_string(img, lang=lang)
        return {"text": text}
    except Exception as e:
        raise HTTPException(500, str(e))


# ============== IMAGE EXTRAS ==============

@router.post("/image/crop")
async def image_crop(file: UploadFile = File(...), left: int = Form(0), top: int = Form(0), right: int = Form(0), bottom: int = Form(0)):
    data = await file.read()
    img = Image.open(io.BytesIO(data))
    w, h = img.size
    box = (left, top, w - right, h - bottom)
    img = img.crop(box)
    out = io.BytesIO(); img.save(out, format="PNG")
    return _stream(out.getvalue(), "cropped.png", "image/png")


@router.post("/image/rotate")
async def image_rotate(file: UploadFile = File(...), angle: float = Form(90)):
    data = await file.read()
    img = Image.open(io.BytesIO(data)).rotate(-angle, expand=True)
    out = io.BytesIO(); fmt = "PNG"; img.save(out, format=fmt)
    return _stream(out.getvalue(), "rotated.png", "image/png")


@router.post("/image/flip")
async def image_flip(file: UploadFile = File(...), direction: str = Form("horizontal")):
    data = await file.read()
    img = Image.open(io.BytesIO(data))
    img = ImageOps.mirror(img) if direction == "horizontal" else ImageOps.flip(img)
    out = io.BytesIO(); img.save(out, format="PNG")
    return _stream(out.getvalue(), "flipped.png", "image/png")


@router.post("/image/border")
async def image_border(file: UploadFile = File(...), thickness: int = Form(20), color: str = Form("#000000")):
    data = await file.read()
    img = Image.open(io.BytesIO(data)).convert("RGB")
    img = ImageOps.expand(img, border=thickness, fill=color)
    out = io.BytesIO(); img.save(out, format="PNG")
    return _stream(out.getvalue(), "bordered.png", "image/png")


@router.post("/image/round-corners")
async def image_round_corners(file: UploadFile = File(...), radius: int = Form(40)):
    data = await file.read()
    img = Image.open(io.BytesIO(data)).convert("RGBA")
    mask = Image.new("L", img.size, 0)
    draw = ImageDraw.Draw(mask)
    draw.rounded_rectangle((0, 0, img.size[0], img.size[1]), radius=radius, fill=255)
    img.putalpha(mask)
    out = io.BytesIO(); img.save(out, format="PNG")
    return _stream(out.getvalue(), "rounded.png", "image/png")


@router.post("/image/dpi")
async def image_change_dpi(file: UploadFile = File(...), dpi: int = Form(300)):
    data = await file.read()
    img = Image.open(io.BytesIO(data))
    out = io.BytesIO()
    fmt = "PNG" if img.mode == "RGBA" else "JPEG"
    img.save(out, format=fmt, dpi=(dpi, dpi))
    return _stream(out.getvalue(), f"dpi-{dpi}.{fmt.lower()}", f"image/{fmt.lower()}")


@router.post("/image/metadata")
async def image_metadata(file: UploadFile = File(...)):
    data = await file.read()
    img = Image.open(io.BytesIO(data))
    info = {"format": img.format, "mode": img.mode, "size": img.size, "info": {k: str(v) for k, v in (img.info or {}).items()}}
    return info


@router.post("/image/remove-metadata")
async def image_remove_metadata(file: UploadFile = File(...)):
    data = await file.read()
    img = Image.open(io.BytesIO(data))
    out_img = Image.new(img.mode, img.size); out_img.putdata(list(img.getdata()))
    out = io.BytesIO(); out_img.save(out, format=img.format or "PNG")
    return _stream(out.getvalue(), "no-metadata." + (img.format or "PNG").lower(), f"image/{(img.format or 'png').lower()}")


@router.post("/image/contact-sheet")
async def image_contact_sheet(files: List[UploadFile] = File(...), cols: int = Form(3), thumb: int = Form(220)):
    imgs = []
    for f in files:
        data = await f.read()
        i = Image.open(io.BytesIO(data)).convert("RGB")
        i.thumbnail((thumb, thumb))
        imgs.append(i)
    if not imgs:
        raise HTTPException(400, "No images")
    rows = (len(imgs) + cols - 1) // cols
    pad = 10
    sheet = Image.new("RGB", (cols * (thumb + pad) + pad, rows * (thumb + pad) + pad), (24, 24, 27))
    for idx, im in enumerate(imgs):
        r, c = divmod(idx, cols)
        sheet.paste(im, (pad + c * (thumb + pad), pad + r * (thumb + pad)))
    out = io.BytesIO(); sheet.save(out, format="JPEG", quality=88)
    return _stream(out.getvalue(), "contact-sheet.jpg", "image/jpeg")


@router.post("/image/favicon")
async def image_favicon(file: UploadFile = File(...)):
    data = await file.read()
    img = Image.open(io.BytesIO(data)).convert("RGBA")
    zbuf = io.BytesIO()
    with zipfile.ZipFile(zbuf, "w", zipfile.ZIP_DEFLATED) as zf:
        for s in [16, 32, 48, 64, 128, 192, 256]:
            r = img.resize((s, s), Image.LANCZOS)
            b = io.BytesIO(); r.save(b, format="PNG")
            zf.writestr(f"favicon-{s}.png", b.getvalue())
        ico = io.BytesIO(); img.resize((64, 64)).save(ico, format="ICO", sizes=[(16, 16), (32, 32), (48, 48), (64, 64)])
        zf.writestr("favicon.ico", ico.getvalue())
    zbuf.seek(0)
    return _stream(zbuf.read(), "favicons.zip", "application/zip")


# ============== BARCODE ==============

@router.post("/barcode/generate")
async def barcode_gen(data: str = Form(...), kind: str = Form("code128")):
    try:
        BarcodeClass = barcode.get_barcode_class(kind)
        bc = BarcodeClass(data, writer=ImageWriter())
        out = io.BytesIO()
        bc.write(out, options={"write_text": True})
        return _stream(out.getvalue(), "barcode.png", "image/png")
    except Exception as e:
        raise HTTPException(500, f"Barcode failed: {e}")


# ============== ARCHIVES ==============

@router.post("/zip/create")
async def zip_create(files: List[UploadFile] = File(...)):
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for f in files:
            zf.writestr(f.filename, await f.read())
    buf.seek(0)
    return _stream(buf.read(), "archive.zip", "application/zip")


@router.post("/zip/extract")
async def zip_extract(file: UploadFile = File(...)):
    data = await file.read()
    try:
        zin = zipfile.ZipFile(io.BytesIO(data))
        names = zin.namelist()
        zin.close()
        return {"files": names, "count": len(names)}
    except Exception as e:
        raise HTTPException(500, str(e))


@router.post("/tar/create")
async def tar_create(files: List[UploadFile] = File(...)):
    buf = io.BytesIO()
    with tarfile.open(fileobj=buf, mode="w") as tf:
        for f in files:
            data = await f.read()
            info = tarfile.TarInfo(name=f.filename); info.size = len(data)
            tf.addfile(info, io.BytesIO(data))
    buf.seek(0)
    return _stream(buf.read(), "archive.tar", "application/x-tar")


@router.post("/gzip/compress")
async def gzip_compress(file: UploadFile = File(...)):
    data = await file.read()
    out = gzip.compress(data)
    return _stream(out, file.filename + ".gz", "application/gzip")


@router.post("/gzip/decompress")
async def gzip_decompress(file: UploadFile = File(...)):
    data = await file.read()
    try:
        out = gzip.decompress(data)
        name = (file.filename or "out").replace(".gz", "")
        return _stream(out, name, "application/octet-stream")
    except Exception as e:
        raise HTTPException(500, str(e))
